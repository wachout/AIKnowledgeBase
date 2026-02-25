import re
from typing import TypedDict, Optional, Dict, Any, Union
from pathlib import Path

from langchain_core.output_parsers import PydanticOutputParser, StrOutputParser
from pydantic import BaseModel, Field
from langgraph.graph import StateGraph, END

# File format parsers - optional imports
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import fitz  # PyMuPDF
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    import PyPDF2
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False

from .prompts import (
    comprehensive_file_analysis_prompt,
)

from dotenv import load_dotenv
from Config.llm_config import get_chat_tongyi

load_dotenv()

# 使用中央化的 LLM 配置获取 ChatTongyi 实例
# 确保在 .env 文件中有正确的 LLM_MODEL_ID、LLM_API_KEY、LLM_BASE_URL 配置
# 注意：如果需要使用特定模型如 qwen-plus，请在 .env 文件中设置 LLM_MODEL_ID=qwen-plus
# temperature=0.3 用于更一致的分析结果
# enable_thinking=False 用于禁用思考过程，有助于结构化输出生成
llm = get_chat_tongyi(temperature=0.3, enable_thinking=False)

# Constants
MAX_CONTENT_LENGTH = 5000  # Maximum characters to analyze


# Pydantic schemas for structured output
class FileMetadata(BaseModel):
    """Schema for file metadata information."""
    file_name: str = Field(description="Name of the analyzed file")
    file_extension: str = Field(description="File extension (e.g., .txt, .md, .py)")
    file_size: int = Field(description="File size in bytes")
    content_length: int = Field(description="Actual content length analyzed (may be truncated)")


class ContentAnalysis(BaseModel):
    """Schema for content analysis results."""
    content_type: str = Field(description="Type of content (text, code, document, etc.)")
    language: Optional[str] = Field(description="Programming language if code file, otherwise None")
    main_topics: list[str] = Field(description="Main topics or themes in the content")
    summary: str = Field(description="Brief summary of the content")


class KeywordAnalysis(BaseModel):
    """Schema for keyword extraction results."""
    keywords: list[str] = Field(description="Important keywords extracted from content")
    key_phrases: list[str] = Field(description="Important phrases or concepts")
    entities: list[str] = Field(description="Named entities (people, organizations, locations)")


class FileAnalysisResult(BaseModel):
    """Complete file analysis result."""
    metadata: FileMetadata
    content_analysis: ContentAnalysis
    keyword_analysis: KeywordAnalysis
    full_analysis: str = Field(description="Detailed analysis text")


# Define the state for our graph
class AgentState(TypedDict):
    """Represents the state of our file analysis agent."""
    input_data: Union[str, Dict[str, Any]]  # Can be file_path string or dict with content
    file_path: Optional[str]  # For file path input
    file_content: Optional[str]  # For direct content input
    metadata: Optional[Dict[str, Any]]  # Raw metadata for description generation
    comprehensive_analysis: Optional[str]  # Complete natural language analysis
    query: Optional[str]  # User's question for targeted analysis
    error: Optional[str]


def parse_file_by_format(file_path: str) -> str:
    """
    Parse file content based on file format.

    Args:
        file_path: Path to the file

    Returns:
        Extracted text content
    """
    path = Path(file_path)
    extension = path.suffix.lower()

    try:
        if extension in ['.txt', '.md', '.py', '.js', '.ts', '.java', '.cpp', '.c', '.cs',
                        '.php', '.rb', '.go', '.rs', '.html', '.css', '.sql', '.sh', '.bash',
                        '.json', '.xml', '.yaml', '.yml']:
            # Plain text files
            try:
                with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            except UnicodeDecodeError:
                with open(path, 'r', encoding='gbk', errors='ignore') as f:
                    return f.read()

        elif extension in ['.docx'] and DOCX_AVAILABLE:
            # Word documents
            doc = DocxDocument(path)
            text = []
            for paragraph in doc.paragraphs:
                text.append(paragraph.text)
            return '\n'.join(text)

        elif extension in ['.pptx', '.ppt'] and PPTX_AVAILABLE:
            # PowerPoint presentations
            prs = Presentation(path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return '\n'.join(text)

        elif extension == '.pdf':
            # PDF documents
            if PDF_AVAILABLE:
                # Try PyMuPDF first
                doc = fitz.open(path)
                text = []
                for page in doc:
                    text.append(page.get_text())
                doc.close()
                return '\n'.join(text)
            elif PYPDF2_AVAILABLE:
                # Fallback to PyPDF2
                with open(path, 'rb') as f:
                    pdf_reader = PyPDF2.PdfReader(f)
                    text = []
                    for page in pdf_reader.pages:
                        text.append(page.extract_text())
                    return '\n'.join(text)
            else:
                raise Exception("PDF parsing libraries not available. Install PyMuPDF or PyPDF2.")

        else:
            # Unsupported format, try as plain text
            try:
                with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            except UnicodeDecodeError:
                with open(path, 'r', encoding='gbk', errors='ignore') as f:
                    return f.read()

    except Exception as e:
        raise Exception(f"Failed to parse {extension} file: {e}")


def read_file_content(file_path: str, max_length: int = MAX_CONTENT_LENGTH) -> str:
    """
    Read file content with length limitation, supporting multiple formats.

    Args:
        file_path: Path to the file
        max_length: Maximum characters to read

    Returns:
        File content as string, truncated if necessary
    """
    try:
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        # Parse content based on file format
        content = parse_file_by_format(file_path)

        # Truncate if too long
        if len(content) > max_length:
            content = content[:max_length]
            print(f"文件内容过长，已截取前 {max_length} 个字符进行分析")

        return content

    except Exception as e:
        raise Exception(f"读取文件失败: {e}")


def process_input_data(input_data: Union[str, Dict[str, Any]]) -> tuple[str, str]:
    """
    Process input data to extract file_path and content.

    Args:
        input_data: Either a file path string or a dict with content info

    Returns:
        Tuple of (file_path, content)
    """
    if isinstance(input_data, str):
        # Input is a file path
        file_path = input_data
        content = read_file_content(file_path)
        return file_path, content
    elif isinstance(input_data, dict):
        # Input is a dict with content
        file_path = input_data.get('file_path', 'unknown_file')
        content = input_data.get('content', '')
        if not content:
            raise ValueError("Content is required when providing dict input")
        return file_path, content
    else:
        raise ValueError("Input data must be a file path string or dict with content")


def extract_file_metadata(file_path: str, content: str) -> Dict[str, Any]:
    """
    Extract basic metadata from file.

    Args:
        file_path: Path to the file
        content: File content

    Returns:
        Metadata dictionary
    """
    path = Path(file_path)

    # Detect programming language or document type based on extension
    language = None
    extension = path.suffix.lower()
    lang_map = {
        # Programming languages
        '.py': 'Python',
        '.js': 'JavaScript',
        '.ts': 'TypeScript',
        '.java': 'Java',
        '.cpp': 'C++',
        '.c': 'C',
        '.cs': 'C#',
        '.php': 'PHP',
        '.rb': 'Ruby',
        '.go': 'Go',
        '.rs': 'Rust',
        '.html': 'HTML',
        '.css': 'CSS',
        '.sql': 'SQL',
        '.sh': 'Shell',
        '.bash': 'Bash',
        # Document formats
        '.md': 'Markdown',
        '.txt': 'Plain Text',
        '.json': 'JSON',
        '.xml': 'XML',
        '.yaml': 'YAML',
        '.yml': 'YAML',
        # Office documents
        '.docx': 'Word Document',
        '.doc': 'Word Document',
        '.pptx': 'PowerPoint Presentation',
        '.ppt': 'PowerPoint Presentation',
        '.pdf': 'PDF Document',
    }
    language = lang_map.get(extension)

    # Get file size (only if file exists)
    file_size = 0
    if path.exists():
        try:
            file_size = path.stat().st_size
        except:
            file_size = 0

    return {
        "file_name": path.name,
        "file_extension": extension,
        "file_size": file_size,
        "content_length": len(content)
    }


# Chain for comprehensive natural language analysis
comprehensive_analysis_chain = comprehensive_file_analysis_prompt | llm | StrOutputParser()


# Define the nodes
def read_file_node(state: AgentState):
    """Node for reading and preprocessing file content."""
    print("---READING FILE---")
    try:
        input_data = state["input_data"]
        file_path, content = process_input_data(input_data)
        metadata = extract_file_metadata(file_path, content)

        # Extract query if provided in input_data
        query = None
        if isinstance(input_data, dict) and "query" in input_data:
            query = input_data["query"]

        result = {
            "file_path": file_path,
            "file_content": content,
            "metadata": metadata
        }

        if query:
            result["query"] = query

        return result
    except Exception as e:
        return {"error": f"Failed to read file: {e}"}


def generate_comprehensive_analysis_node(state: AgentState):
    """Node for generating comprehensive natural language analysis."""
    print("---GENERATING COMPREHENSIVE ANALYSIS---")
    try:
        if state.get("error"):
            print("DEBUG: Node has error, returning empty")
            return {}

        content = state.get("file_content", "")
        metadata = state.get("metadata", {})
        query = state.get("query", "")

        print(f"DEBUG: Content length: {len(content)}")
        print(f"DEBUG: Metadata: {metadata}")
        print(f"DEBUG: Query: {query}")

        if not content:
            print("DEBUG: No content available")
            return {"error": "No content available for analysis"}

        print("DEBUG: Invoking comprehensive analysis chain...")
        # Generate comprehensive natural language analysis using LLM
        # Include query in the analysis if provided
        analysis_input = {
            "file_name": metadata.get("file_name", "unknown"),
            "file_extension": metadata.get("file_extension", ""),
            "file_size": metadata.get("file_size", 0),
            "content_length": metadata.get("content_length", 0),
            "content": content,
            "query": query if query else "请对这个文件进行全面分析"  # Provide default query if none provided
        }

        result = comprehensive_analysis_chain.invoke(analysis_input)

        print(f"DEBUG: LLM result type: {type(result)}")
        print(f"DEBUG: LLM result length: {len(result) if result else 0}")
        print(f"DEBUG: LLM result preview: {result[:100] if result else 'None'}...")

        cleaned_result = result.strip() if result else ""
        print(f"DEBUG: Returning comprehensive_analysis: {len(cleaned_result)} chars")

        return {"comprehensive_analysis": cleaned_result}
    except Exception as e:
        print(f"DEBUG: Exception in comprehensive analysis: {e}")
        import traceback
        traceback.print_exc()
        return {"error": f"Failed to generate comprehensive analysis: {e}"}


def run_file_analysis(input_data: Union[str, Dict[str, Any]]) -> Dict[str, Any]:
    """
    Run the file analysis agent.

    Args:
        input_data: Either a file path string or dict with content info

    Returns:
        Analysis result dictionary
    """
    # Define the graph
    workflow = StateGraph(AgentState)

    # Add nodes
    workflow.add_node("read_file", read_file_node)
    workflow.add_node("generate_comprehensive_analysis", generate_comprehensive_analysis_node)

    # Define edges - simple sequential flow
    workflow.set_entry_point("read_file")
    workflow.add_edge("read_file", "generate_comprehensive_analysis")
    workflow.add_edge("generate_comprehensive_analysis", END)

    # Compile the graph
    app = workflow.compile()

    # Run the analysis
    initial_state = {
        "input_data": input_data,
        "file_path": None,
        "file_content": None,
        "metadata": None,
        "comprehensive_analysis": None,
        "query": None,
        "error": None
    }

    try:
        result = app.invoke(initial_state)

        # 调试：打印完整结果状态
        print(f"DEBUG: Final result state keys: {list(result.keys())}")
        print(f"DEBUG: comprehensive_analysis present: {'comprehensive_analysis' in result}")
        print(f"DEBUG: comprehensive_analysis value: {result.get('comprehensive_analysis', 'NOT_FOUND')}")

        if result.get("error"):
            file_path_result = result.get("file_path", "unknown_file")
            return {
                "success": False,
                "error": result["error"],
                "file_path": file_path_result
            }

        # 返回综合分析结果
        comprehensive_analysis = result.get("comprehensive_analysis", "")
        file_path_result = result.get("file_path", "unknown_file")

        if not comprehensive_analysis:
            return {
                "success": False,
                "error": "Analysis completed but no comprehensive analysis was generated",
                "file_path": file_path_result
            }

        return {
            "success": True,
            "result": comprehensive_analysis,  # 返回完整的自然语言分析
            "file_path": file_path_result
        }

    except Exception as e:
        file_path_result = "unknown_file"
        try:
            if isinstance(input_data, str):
                file_path_result = input_data
            elif isinstance(input_data, dict):
                file_path_result = input_data.get("file_path", "unknown_file")
        except:
            pass

        return {
            "success": False,
            "error": str(e),
            "file_path": file_path_result
        }


# Convenience functions for easy usage
def analyze_file(file_path: str) -> str:
    """
    Analyze a file and return natural language analysis.

    Args:
        file_path: Path to the file to analyze

    Returns:
        Natural language analysis text
    """
    result = run_file_analysis(file_path)

    if not result["success"]:
        raise Exception(f"File analysis failed: {result['error']}")

    if result["result"] is None:
        raise Exception("File analysis completed but no result was generated")

    return result["result"]


def analyze_content(content: str, file_name: str = "unknown_file") -> str:
    """
    Analyze content directly and return natural language analysis.

    Args:
        content: Text content to analyze
        file_name: Optional file name for metadata

    Returns:
        Natural language analysis text
    """
    input_data = {
        "file_path": file_name,
        "content": content
    }

    result = run_file_analysis(input_data)

    if not result["success"]:
        raise Exception(f"Content analysis failed: {result['error']}")

    if result["result"] is None:
        raise Exception("Content analysis completed but no result was generated")

    return result["result"]


def get_supported_formats() -> Dict[str, str]:
    """
    Get information about supported file formats and required libraries.

    Returns:
        Dictionary mapping file extensions to format descriptions
    """
    formats = {
        # Plain text formats (no external libraries needed)
        '.txt': 'Plain text files',
        '.md': 'Markdown files',
        '.py': 'Python source code',
        '.js': 'JavaScript source code',
        '.ts': 'TypeScript source code',
        '.java': 'Java source code',
        '.cpp': 'C++ source code',
        '.c': 'C source code',
        '.cs': 'C# source code',
        '.php': 'PHP source code',
        '.rb': 'Ruby source code',
        '.go': 'Go source code',
        '.rs': 'Rust source code',
        '.html': 'HTML files',
        '.css': 'CSS files',
        '.sql': 'SQL files',
        '.sh': 'Shell scripts',
        '.bash': 'Bash scripts',
        '.json': 'JSON files',
        '.xml': 'XML files',
        '.yaml': 'YAML files',
        '.yml': 'YAML files',
    }

    # Conditionally add formats based on available libraries
    if DOCX_AVAILABLE:
        formats['.docx'] = 'Microsoft Word documents (python-docx required)'
    else:
        formats['.docx'] = 'Microsoft Word documents (python-docx NOT installed)'

    if PPTX_AVAILABLE:
        formats['.pptx'] = 'PowerPoint presentations (python-pptx required)'
        formats['.ppt'] = 'PowerPoint presentations (python-pptx required)'
    else:
        formats['.pptx'] = 'PowerPoint presentations (python-pptx NOT installed)'
        formats['.ppt'] = 'PowerPoint presentations (python-pptx NOT installed)'

    if PDF_AVAILABLE or PYPDF2_AVAILABLE:
        library = 'PyMuPDF' if PDF_AVAILABLE else 'PyPDF2'
        formats['.pdf'] = f'PDF documents ({library} required)'
    else:
        formats['.pdf'] = 'PDF documents (PyMuPDF or PyPDF2 NOT installed)'

    return formats
